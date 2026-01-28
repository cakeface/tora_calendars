#!/usr/bin/env python3
"""
Generate a 2026 calendar PowerPoint with Monday-Sunday week layout.
Saturday and Sunday appear on the far right of each calendar.
"""

import calendar
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml

# Set Monday as first day of week (0 = Monday, 6 = Sunday)
calendar.setfirstweekday(calendar.MONDAY)

YEAR = 2026
MONTH_NAMES = [
    "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December"
]
DAY_NAMES = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]

# Colors (hex strings for XML)
HEADER_BG_HEX = "2C3E50"  # Dark blue
WEEKEND_BG_HEX = "ECF0F1"  # Light gray for weekend columns
WHITE_HEX = "FFFFFF"

# Colors (RGB for text)
HEADER_BG = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)


def set_cell_fill(cell, hex_color):
    """Set the background fill color of a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing fill
    for child in list(tcPr):
        if 'solidFill' in child.tag:
            tcPr.remove(child)
    # Add new solid fill
    solidFill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{hex_color}"/></a:solidFill>'
    )
    tcPr.append(solidFill)


def set_cell_borders(cell, hex_color, width_pt=1):
    """Set all four borders of a table cell to a solid color."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Width in EMUs (English Metric Units) - 12700 EMUs per point
    width_emu = int(width_pt * 12700)

    # Border element names for all four sides
    border_names = ['lnL', 'lnR', 'lnT', 'lnB']  # left, right, top, bottom

    # Remove existing border elements
    for child in list(tcPr):
        for name in border_names:
            if name in child.tag:
                tcPr.remove(child)
                break

    # Add borders for all four sides
    for border_name in border_names:
        border_xml = parse_xml(
            f'<a:{border_name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width_emu}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'<a:prstDash val="solid"/>'
            f'</a:{border_name}>'
        )
        tcPr.append(border_xml)


def create_month_slide(prs, month):
    """Create a slide for a single month."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Get the calendar for this month
    cal = calendar.monthcalendar(YEAR, month)

    # Add month/year title - compact, left-aligned
    title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(5), Inches(0.4))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"{MONTH_NAMES[month-1]} {YEAR}"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = HEADER_BG
    title_para.alignment = PP_ALIGN.LEFT

    # Table dimensions
    rows = len(cal) + 1  # +1 for header row
    cols = 7

    # Table position and size - maximize space, minimal margins
    left = Inches(0.2)
    top = Inches(0.6)
    width = Inches(9.6)
    height = Inches(6.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths (equal)
    col_width = Inches(9.6 / 7)
    for i in range(cols):
        table.columns[i].width = col_width

    # Set row heights - small header, large day rows
    header_height = Inches(0.35)
    day_row_height = Inches((6.8 - 0.35) / len(cal))
    table.rows[0].height = header_height
    for i in range(1, rows):
        table.rows[i].height = day_row_height

    # Header row with day names - compact
    for i, day_name in enumerate(DAY_NAMES):
        cell = table.cell(0, i)
        cell.text = day_name
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_cell_fill(cell, HEADER_BG_HEX)
        set_cell_borders(cell, HEADER_BG_HEX)

    # Fill in the days - number in top-right corner
    for week_idx, week in enumerate(cal):
        row_idx = week_idx + 1
        for day_idx, day in enumerate(week):
            cell = table.cell(row_idx, day_idx)
            if day != 0:
                cell.text = str(day)
            else:
                cell.text = ""

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = BLACK
            para.alignment = PP_ALIGN.RIGHT  # Number in top-right
            cell.vertical_anchor = MSO_ANCHOR.TOP  # Number at top

            # Weekend columns (Saturday=5, Sunday=6) get light gray background
            if day_idx >= 5:
                set_cell_fill(cell, WEEKEND_BG_HEX)
            else:
                set_cell_fill(cell, WHITE_HEX)

            # Dark grid borders
            set_cell_borders(cell, HEADER_BG_HEX)

    return slide


def validate_calendar(prs):
    """Validate that all 12 months are correctly rendered."""
    issues = []

    for month in range(1, 13):
        slide = prs.slides[month - 1]

        # Get the table from the slide
        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break

        if table is None:
            issues.append(f"{MONTH_NAMES[month-1]}: No table found")
            continue

        # Get expected calendar data
        cal = calendar.monthcalendar(YEAR, month)
        expected_days = [day for week in cal for day in week if day != 0]

        # Extract days from table
        found_days = []
        for row_idx in range(1, len(cal) + 1):
            for col_idx in range(7):
                cell_text = table.cell(row_idx, col_idx).text.strip()
                if cell_text:
                    found_days.append(int(cell_text))

        # Check day count
        days_in_month = calendar.monthrange(YEAR, month)[1]
        if len(found_days) != days_in_month:
            issues.append(f"{MONTH_NAMES[month-1]}: Expected {days_in_month} days, found {len(found_days)}")

        # Check that all days are present
        expected_set = set(range(1, days_in_month + 1))
        found_set = set(found_days)
        if expected_set != found_set:
            missing = expected_set - found_set
            extra = found_set - expected_set
            if missing:
                issues.append(f"{MONTH_NAMES[month-1]}: Missing days {missing}")
            if extra:
                issues.append(f"{MONTH_NAMES[month-1]}: Extra days {extra}")

        # Verify first day of month is in correct column
        first_day_weekday = date(YEAR, month, 1).weekday()  # 0=Monday, 6=Sunday
        # Find where day 1 appears
        for col_idx in range(7):
            cell_text = table.cell(1, col_idx).text.strip()
            if cell_text == "1":
                if col_idx != first_day_weekday:
                    issues.append(f"{MONTH_NAMES[month-1]}: Day 1 in wrong column. Expected {DAY_NAMES[first_day_weekday]}, found {DAY_NAMES[col_idx]}")
                break

    return issues


def main():
    print(f"Creating 2026 Calendar (Monday-Sunday layout)...")

    # Create presentation with widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create a slide for each month
    for month in range(1, 13):
        create_month_slide(prs, month)
        print(f"  Created {MONTH_NAMES[month-1]}")

    # Validate the calendar
    print("\nValidating calendar...")
    issues = validate_calendar(prs)

    if issues:
        print("VALIDATION ERRORS:")
        for issue in issues:
            print(f"  - {issue}")
    else:
        print("All 12 months validated successfully!")

        # Print summary of each month
        print("\nCalendar Summary:")
        for month in range(1, 13):
            days_in_month = calendar.monthrange(YEAR, month)[1]
            first_day = date(YEAR, month, 1).weekday()
            print(f"  {MONTH_NAMES[month-1]:12} - {days_in_month} days, starts on {DAY_NAMES[first_day]}")

    # Save the presentation
    output_file = "2026_Calendar.pptx"
    prs.save(output_file)
    print(f"\nCalendar saved to: {output_file}")

    return len(issues) == 0


if __name__ == "__main__":
    success = main()
    exit(0 if success else 1)
